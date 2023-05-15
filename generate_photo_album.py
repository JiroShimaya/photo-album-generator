import os
import re

import random
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def generate_photo_album(input_dir, rows=2, columns=3, num_pages=5, crop_aspect_ratio=1.0, output_file="photo_album.pptx"):
    # パワーポイントのプレゼンテーションを作成
    prs = Presentation()

    # A4横のサイズを設定
    slide_width = Inches(11.69)  # 29.7 cm
    slide_height = Inches(8.27)  # 21 cm
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    # 写真のリストを取得
    photo_list = get_photo_list(input_dir)
    # 最大写真数
    max_num_photos = rows * columns * num_pages
    # 写真をランダムに取得
    photo_list = random.sample(photo_list, min(len(photo_list), max_num_photos))
    # 写真を撮影日時の昇順でソート
    photo_list = sort_photos_by_datetime(photo_list)

    # ページごとに写真を配置
    for page in range(num_pages):
        slide_layout = prs.slide_layouts[6]  # 空のスライドレイアウトを選択
        slide = prs.slides.add_slide(slide_layout)

        # 必要な数の写真を先頭から取得してリストを作成
        num_photos_needed = rows * columns
        selected_photos = photo_list[page*num_photos_needed:min(len(photo_list), (page+1)*num_photos_needed)]

        # 写真をグリッド状に配置（画像をトリミング）
        add_photos_to_slide_grid(slide, slide_width, slide_height, rows, columns, selected_photos, crop_aspect_ratio)

        # 写真がなくなったら修了
        if len(photo_list) <= num_photos_needed * (page+1):
            break


    # プレゼンテーションを保存
    prs.save(output_file)
    print("写真アルバムが生成されました。")


def sort_photos_by_datetime(file_paths):
    # 撮影日時の昇順でソート
    sorted_file_paths = sorted(file_paths, key=get_datetime_taken)

    return sorted_file_paths

def get_datetime_taken(file_path):
    try:
        with Image.open(file_path) as img:
            exif_data = img._getexif()
            datetime_taken = exif_data.get(36867)  # Exifタグのキー(36867)は撮影日時を表す
            if datetime_taken:
                return datetime_taken
    except (AttributeError, KeyError, IndexError):
        pass

    # 撮影日時が取得できなかった場合はファイル名の昇順でソート
    return os.path.basename(file_path)

def is_horizontal_image(image_path):
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            return width > height
    except (OSError, FileNotFoundError):
        return False
    
def get_photo_list(input_dir):
    # フォルダ内の写真のパスをリストとして取得
    photo_list = []
    for file_name in os.listdir(input_dir):
        if file_name.endswith(".jpg") or file_name.endswith(".jpeg") or file_name.endswith(".png"):
            photo_list.append(os.path.join(input_dir, file_name))
    # 撮影日の昇順に写真をソート
    photo_list = sort_photos_by_datetime(photo_list)
    # 横長の写真のみを選択
    photo_list = [photo_path for photo_path in photo_list if is_horizontal_image(photo_path)]
    return photo_list

def crop_image_with_aspect_ratio(image_path, aspect_ratio, output_path):
    # 画像を開く
    image = Image.open(image_path)

    # 画像のサイズとアスペクト比を取得
    width, height = image.size
    image_aspect_ratio = width / height

    if image_aspect_ratio > aspect_ratio:
        # 画像の幅がアスペクト比より大きい場合、横方向にトリミング
        new_width = int(height * aspect_ratio)
        left = (width - new_width) // 2
        right = left + new_width
        image = image.crop((left, 0, right, height))
    else:
        # 画像の高さがアスペクト比より大きい場合、縦方向にトリミング
        new_height = int(width / aspect_ratio)
        top = (height - new_height) // 2
        bottom = top + new_height
        image = image.crop((0, top, width, bottom))

    # トリミングされた画像を保存
    image.save(output_path)

def add_photos_to_slide_grid(slide, slide_width, slide_height, rows, columns, photo_list, crop_aspect_ratio, 
                             margin_outer_top=0.05, margin_outer_bottom=0.05, margin_outer_left=0.05, margin_outer_right=0.05, 
                             margin_inner_top=0.05, margin_inner_bottom=0.05, margin_inner_left=0.05, margin_inner_right=0.05):
    # 外側マージンを計算
    outer_margin_top = slide_height * margin_outer_top
    outer_margin_bottom = slide_height * margin_outer_bottom
    outer_margin_left = slide_width * margin_outer_left
    outer_margin_right = slide_width * margin_outer_right

    # グリッドセルの幅と高さを計算（外側マージンを除く）
    grid_width = (slide_width - outer_margin_left - outer_margin_right) / columns
    grid_height = (slide_height - outer_margin_top - outer_margin_bottom) / rows

    # 内側マージン（写真間の隙間）を計算
    margin_top = grid_height * margin_inner_top
    margin_bottom = grid_height * margin_inner_bottom
    margin_left = grid_width * margin_inner_left
    margin_right = grid_width * margin_inner_right

    # 写真をグリッド状に配置
    for row in range(rows):
        for col in range(columns):
            if not photo_list:  # 全ての写真が選ばれたら終了
                break

            photo_path = photo_list.pop(0)  # 写真を選択してリストから削除

            # 画像をトリミングして指定されたアスペクト比に合わせる
            output_path = "temp.jpg"  # 一時的なファイルに保存
            crop_image_with_aspect_ratio(photo_path, crop_aspect_ratio, output_path)

            # クロップされた画像のアスペクト比を維持したまま、グリッド内にできるだけ大きく表示
            photo_width = grid_width - margin_left - margin_right
            photo_height = photo_width / crop_aspect_ratio

            if photo_height > grid_height - margin_top - margin_bottom:
                photo_height = grid_height - margin_top - margin_bottom
                photo_width = photo_height * crop_aspect_ratio

            # 写真をスライドに貼り付ける
            left = outer_margin_left + col * grid_width + (grid_width - photo_width) / 2
            top = outer_margin_top + row * grid_height + (grid_height - photo_height) / 2
            slide.shapes.add_picture(output_path, left, top, width=photo_width, height=photo_height)

            # 一時的なファイルを削除
            os.remove(output_path)

def parse_aspect_ratio(value):
    match = re.match(r"(\d+):(\d+)", value)
    if match:
        numerator = int(match.group(1))
        denominator = int(match.group(2))
        return numerator / denominator
    else:
        return float(value)

if __name__ == '__main__':
    # コマンドライン引数の解析
    parser = argparse.ArgumentParser(description='Generate a photo album in PowerPoint.')
    parser.add_argument('--input_dir', type=str, help='Input directory path')
    parser.add_argument('--rows', type=int, default=2, help='Number of rows of photos per page')
    parser.add_argument('--columns', type=int, default=3, help='Number of columns of photos per page')
    parser.add_argument('--num_pages', type=int, default=5, help='Total number of pages')
    parser.add_argument('--crop_aspect_ratio', type=parse_aspect_ratio, default=1.0, help='Aspect ratio for photo cropping (e.g., 16:9)')
    parser.add_argument('--output_file', type=str, default='photo_album.pptx', help='Output file name')

    args = parser.parse_args()

    # 使用例
    input_dir = args.input_dir
    rows = args.rows
    columns = args.columns
    num_pages = args.num_pages
    crop_aspect_ratio = args.crop_aspect_ratio
    output_file = args.output_file

    generate_photo_album(input_dir, rows, columns, num_pages, crop_aspect_ratio, output_file)
